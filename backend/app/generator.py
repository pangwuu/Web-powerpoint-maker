import os
from random import choice
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from deep_translator import GoogleTranslator
from functools import cache
import io

# Import our models and helpers
from .models import Song, SongSection, GenerateRequest, AnnouncementItem, OfferingInfo
from .ai_translate import translate_with_gemini, translate_text_gemini
from .bible import get_correct_copyright_message

# Paths
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
BACKEND_DIR = os.path.dirname(BASE_DIR)
TEMPLATES_DIR = os.path.join(BACKEND_DIR, 'templates')
ASSETS_DIR = os.path.join(BACKEND_DIR, 'assets')

def get_template_path(size: str = "medium") -> str:
    """Finds a template file in the templates directory based on size."""
    if not os.path.exists(TEMPLATES_DIR):
        raise FileNotFoundError(f"Templates directory not found: {TEMPLATES_DIR}")
    
    all_templates = []
    for root, dirs, files in os.walk(TEMPLATES_DIR):
        for file in files:
            if file.endswith(".pptx"):
                all_templates.append(os.path.join(root, file))
    
    if not all_templates:
        raise FileNotFoundError("No .pptx templates found.")
    
    # Filter by size if possible
    size_files = [f for f in all_templates if size.lower() in f.lower()]
    if size_files:
        return choice(size_files)
    
    return choice(all_templates)

def create_blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank slide layout
    return prs.slides.add_slide(layout)

def add_text_to_slide(slide, text, prs, font_size, alignment=PP_ALIGN.CENTER, position_percent=0.35,
                     colour=None, bold=False, italic=False, underline=False):
    width = prs.slide_width
    height = prs.slide_height
    left = int((width - width * 0.8) / 2)
    top = int(height * position_percent)
    text_box = slide.shapes.add_textbox(left, top, int(width * 0.8), int(height * 0.5))
    text_frame = text_box.text_frame
    text_frame.text = text
    
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.alignment = alignment
        paragraph.font.bold = bold
        paragraph.font.italic = italic
        paragraph.font.underline = underline
        if colour:
            try:
                paragraph.font.color.rgb = RGBColor.from_string(colour)
            except:
                pass # Fallback
    
    text_frame.word_wrap = True

def create_title_slide(title_text: str, subtitle_text: str, prs, title_size, default_body_size=8):
    blank_slide = create_blank_slide(prs)
    subtitle_text = subtitle_text.replace("\n", " ")
    add_text_to_slide(blank_slide, title_text, prs, title_size, position_percent=0.2)
    add_text_to_slide(blank_slide, subtitle_text, prs, default_body_size, position_percent=0.6)
    return prs

def create_text_slide(body_text, prs, body_size, slide_number=0, total_slides=0):
    blank_slide_layout = prs.slide_layouts[6]
    lyric_slide = prs.slides.add_slide(blank_slide_layout)

    body_width = prs.slide_width * 0.9
    body_height = prs.slide_height * 0.8
    body_left = (prs.slide_width - body_width) / 2
    body_top = prs.slide_height * 0.1

    body_box = lyric_slide.shapes.add_textbox(left=body_left, top=body_top, width=body_width, height=body_height)
    body_frame = body_box.text_frame
    body_frame.text = body_text
    body_frame.word_wrap = True

    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(body_size)
        paragraph.alignment = PP_ALIGN.CENTER
    
    return prs

def create_title_and_text_slide(title_text, body_text, prs, title_size, body_size):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    title_width = prs.slide_width * 0.9
    title_height = prs.slide_height * 0.15
    title_left = (prs.slide_width - title_width) / 2
    title_top = Inches(0.3)

    title_box = slide.shapes.add_textbox(left=title_left, top=title_top, width=title_width, height=title_height)
    title_frame = title_box.text_frame
    title_frame.text = title_text

    body_width = prs.slide_width * 0.9
    body_height = prs.slide_height * 0.8
    body_left = (prs.slide_width - body_width) / 2
    body_top = title_height + title_top

    body_box = slide.shapes.add_textbox(left=body_left, top=body_top, width=body_width, height=body_height)
    body_frame = body_box.text_frame
    body_frame.text = body_text

    title_frame.word_wrap = True
    body_frame.word_wrap = True

    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(title_size)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(body_size)
        paragraph.alignment = PP_ALIGN.CENTER
    
    return prs

def add_title_with_image_on_right(prs: Presentation, title_text: str, image_type: str, left_text_size: int) -> Presentation:
    margin_right = margin_top = prs.slide_height * 0.1
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_title = slide.shapes.add_textbox(prs.slide_width * 0.05, margin_top, prs.slide_width * 0.4, prs.slide_height - 2 * margin_top)
    title_text_frame = slide_title.text_frame
    title_text_frame.word_wrap = True
    p = title_text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(left_text_size)

    directory_path = os.path.join(ASSETS_DIR, image_type.strip().title())
    if not os.path.exists(directory_path):
        return prs # Skip if not found

    files = [f for f in os.listdir(directory_path) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    if not files:
        return prs

    image_path = os.path.join(directory_path, choice(files))
    image_width = image_height = prs.slide_height - 2 * margin_top
    
    try:
        slide.shapes.add_picture(image_path, prs.slide_width - image_width - margin_right, margin_top, width=image_width, height=image_height)
    except Exception as e:
        print(f"Error adding image: {e}")

    return prs

def create_bulletin_slide(slide, prs, date, songs: List[str], verses: List[str], response_songs: List[str], speaker: str, topic: str, church_name: str, service_name: str):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)

    ppt_text_break = "\x0b"
    # Format date nicely
    formatted_date = date.replace("-", ".")
    
    add_text_to_slide(slide, church_name, prs, 30, position_percent=0.05, bold=True, colour="000000")
    add_text_to_slide(slide, f"Welcome to Our {service_name}{ppt_text_break}{formatted_date}", prs, 23, position_percent=0.15, bold=True, colour="000000")
    add_text_to_slide(slide, f"{service_name} 9:45 – 11:00 am", prs, 20, position_percent=0.35, colour="7030A0", bold=True, underline=True)

    bulletin_summary = {
        "left": [
            f"Worship Songs:{ppt_text_break * (len(songs) - 1)}",
            f"Passage:{ppt_text_break * (len(verses) - 1)}",
            "Speaker:",
            "Topic:",
            "Response Song:",
            "Announcements and Closing Prayer",
            "Benediction",
        ],
        "right": [
            ppt_text_break.join(songs),
            ppt_text_break.join(verses),
            speaker,
            topic,
            ppt_text_break.join(response_songs),
            "",
            "",
        ],
    }

    add_text_to_slide(slide, ppt_text_break.join(bulletin_summary["left"]), prs, 20, position_percent=0.4, alignment=PP_ALIGN.LEFT, bold=True, colour="000000")
    add_text_to_slide(slide, ppt_text_break.join(bulletin_summary["right"]), prs, 20, position_percent=0.4, alignment=PP_ALIGN.CENTER, colour="000000")

def create_offering_slide(prs: Presentation, tithing_heading_size: int, tithing_body_size: int, offering_info: 'OfferingInfo'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_height = Inches(1)
    title_top = (prs.slide_height - title_height) / 2
    left_textbox = prs.slide_width*0.075
    top_textbox = title_top
    width_textbox = prs.slide_width*0.3
    height_textbox = Inches(1)

    textbox = slide.shapes.add_textbox(left_textbox, top_textbox - height_textbox/2, width_textbox, height_textbox)
    p = textbox.text_frame.add_paragraph()
    p.text = 'Offering'
    p.font.size = Pt(tithing_heading_size)

    # Boxes
    box_height = prs.slide_height/6
    box_spacing = box_height*0.15
    total_height = 5 * box_height + 4 * box_spacing
    top_first_box = (prs.slide_height - total_height) / 2

    # Data for offering
    offering_data = [
        (f"Account name: {offering_info.account_name}", "church.png"),
        (f"Account number: {offering_info.account_number}", "account_number.png"),
        (f"BSB: {offering_info.bsb}", "bsb.png"),
        (f"Please put in \"{offering_info.reference}\" as the reference", "hands.png"),
        (offering_info.details, "box.png"),
    ]

    for i, (text, img_file) in enumerate(offering_data):
        width_rect = 0.65 * prs.slide_width
        height_rect = box_height
        left_rect = prs.slide_width - (top_first_box + width_rect)
        top_rect = top_first_box + i * (box_spacing + box_height)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_rect, top_rect, width_rect, height_rect)
        fill = shape.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.BACKGROUND_2
        fill.fore_color.brightness = -0.15
        shape.line.fill.background()

        left_textbox = left_rect + width_rect/7
        textbox = slide.shapes.add_textbox(left_textbox, top_rect, width_rect - width_rect/5, height_rect)
        p = textbox.text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(tithing_body_size)
        textbox.text_frame.word_wrap = True

        # Image
        img_path = os.path.join(ASSETS_DIR, 'Tithing', img_file)
        if os.path.exists(img_path):
            image_width = image_height = box_height*0.8
            try:
                slide.shapes.add_picture(img_path, left_rect + box_height*0.1, top_rect + box_height*0.1, image_width, image_height)
            except Exception:
                pass

    return prs

@cache
def translate_text(text: str, language: str = 'Mandarin Chinese') -> str:
    # Try Gemini first
    gemini_translated = translate_text_gemini(text, language)
    if gemini_translated:
        return gemini_translated

    # Fallback to Google Translate
    lang_map = {
        "afrikaans": "af", "albanian": "sq", "amharic": "am", "arabic": "ar", "armenian": "hy", 
        "azerbaijani": "az", "basque": "eu", "belarusian": "be", "bengali": "bn", "bosnian": "bs", 
        "bulgarian": "bg", "catalan": "ca", "cebuano": "ceb", "chichewa": "ny", 
        "chinese (simplified)": "zh-CN", "mandarin chinese": "zh-CN", "chinese (traditional)": "zh-TW", 
        "corsican": "co", "croatian": "hr", "czech": "cs", "danish": "da", "dutch": "nl", 
        "english": "en", "esperanto": "eo", "estonian": "et", "filipino": "tl", "finnish": "fi", 
        "french": "fr", "frisian": "fy", "galician": "gl", "georgian": "ka", "german": "de", 
        "greek": "el", "gujarati": "gu", "haitian creole": "ht", "hausa": "ha", "hawaiian": "haw", 
        "hebrew": "iw", "hindi": "hi", "hmong": "hmn", "hungarian": "hu", "icelandic": "is", 
        "igbo": "ig", "indonesian": "id", "irish": "ga", "italian": "it", "japanese": "ja", 
        "javanese": "jw", "kannada": "kn", "kazakh": "kk", "khmer": "km", "korean": "ko", 
        "kurdish (kurmanji)": "ku", "kyrgyz": "ky", "lao": "lo", "latin": "la", "latvian": "lv", 
        "lithuanian": "lt", "luxembourgish": "lb", "macedonian": "mk", "malagasy": "mg", 
        "malay": "ms", "malayalam": "ml", "maltese": "mt", "maori": "mi", "marathi": "mr", 
        "mongolian": "mn", "myanmar (burmese)": "my", "nepali": "ne", "norwegian": "no", 
        "odia": "or", "pashto": "ps", "persian": "fa", "polish": "pl", "portuguese": "pt", 
        "punjabi": "pa", "romanian": "ro", "russian": "ru", "samoan": "sm", "scots gaelic": "gd", 
        "serbian": "sr", "sesotho": "st", "shona": "sn", "sindhi": "sd", "sinhala": "si", 
        "slovak": "sk", "slovenian": "sl", "somali": "so", "spanish": "es", "sundanese": "su", 
        "swahili": "sw", "swedish": "sv", "tajik": "tg", "tamil": "ta", "telugu": "te", 
        "thai": "th", "turkish": "tr", "ukrainian": "uk", "urdu": "ur", "uyghur": "ug", 
        "uzbek": "uz", "vietnamese": "vi", "welsh": "cy", "xhosa": "xh", "yiddish": "yi", 
        "yoruba": "yo", "zulu": "zu"
    }
    code = lang_map.get(language.lower(), "zh-CN")
    try:
        return GoogleTranslator(source='auto', target=code).translate(text)
    except Exception:
        return text

def append_song(prs, song: Song, title_size, font_size, translate: bool = False, language: str = "Chinese (Simplified)"):
    translation_map = {}
    if translate:
        # Collect all unique lines from the song to translate in one go
        all_lines = []
        for section in song.sections:
            section_lines = [l.strip() for l in section.content.split('\n') if l.strip()]
            all_lines.extend(section_lines)
        
        # Deduplicate while preserving order
        unique_lines = []
        for l in all_lines:
            if l not in unique_lines:
                unique_lines.append(l)
        
        if unique_lines:
            text_to_translate = "\n".join(unique_lines)
            translated_text_block = translate_text(text_to_translate, language)
            translated_lines = translated_text_block.split('\n')
            
            # Map original lines to translated lines
            for i, original in enumerate(unique_lines):
                if i < len(translated_lines):
                    translation_map[original] = translated_lines[i].strip()
                else:
                    translation_map[original] = original # Fallback to original if count mismatch

    # Title Slide
    ccli_info = ""
    if song.ccli_number:
         ccli_info = f"CCLI Licence No. {song.ccli_number}"
    
    if translate:
        create_title_slide_translated(song.title, ccli_info, prs, title_size, 8, language)
    else:
        create_title_slide(song.title, ccli_info, prs, title_size)
    
    # Lyrics Slides
    for section in song.sections:
        original_content = section.content.strip()
        if not original_content:
            continue

        # Split into lines
        lines = [line.strip() for line in original_content.split('\n') if line.strip()]
        
        # Chunk lines: if translating, we use smaller chunks to fit both languages
        chunk_size = 2 if translate else 4
        chunks = [lines[i:i + chunk_size] for i in range(0, len(lines), chunk_size)]
        
        for chunk in chunks:
            if translate:
                # Interleave original and translated lines from the pre-translated map
                combined_lines = []
                for line in chunk:
                    t_line = translation_map.get(line, line)
                    combined_lines.append(line)
                    combined_lines.append(t_line)
                
                final_text = "\n".join(combined_lines)
                create_text_slide(final_text, prs, font_size * 0.85) 
            else:
                final_text = "\n".join(chunk)
                create_text_slide(final_text, prs, font_size)

    return prs

def create_title_slide_translated(title_text, subtitle_text, prs, title_size, subtitle_size, language):
    blank_slide = create_blank_slide(prs)
    t_title = translate_text(title_text, language)
    
    full_title = f"{title_text}\n{t_title}"
    add_text_to_slide(blank_slide, full_title, prs, title_size, position_percent=0.2)
    add_text_to_slide(blank_slide, subtitle_text, prs, subtitle_size, position_percent=0.6)
    return prs

def generate_powerpoint(request: GenerateRequest) -> io.BytesIO:
    template_path = get_template_path(request.template_name)
    prs = Presentation(template_path)

    # Font sizes
    font_map = {
        'small': {'title': 70, 'song': 53, 'bible': 43, 'tithing': 32},
        'medium': {'title': 50, 'song': 33, 'bible': 32, 'tithing': 23},
        'large': {'title': 40, 'song': 27, 'bible': 23, 'tithing': 16}
    }
    # Guess size from template name
    size = 'medium'
    if 'small' in os.path.basename(template_path): size = 'small'
    elif 'large' in os.path.basename(template_path): size = 'large'
    fonts = font_map[size]

    # 1. Start
    create_blank_slide(prs) # Bulletin placeholder (index 0)
    create_title_slide(request.church_name, request.service_name, prs, fonts['title'])

    # 2. Songs
    song_names = [s.title for s in request.songs]
    for song in request.songs:
        append_song(prs, song, fonts['title'], fonts['song'], request.translate, request.language)

    # 3. Communion (Detect first Sunday logic can be done in frontend or here)
    # We'll just assume if user wants it, they add a generic "Communion" slide item, but 
    # for now we follow the script logic: if date is first 7 days.
    try:
        day = int(request.date.split("-")[-1]) # Assuming YYYY-MM-DD
        if 1 <= day <= 7:
             add_title_with_image_on_right(prs, "Holy Communion", "Communion", fonts['title'] - 10)
    except:
        pass

    # 4. Bible
    add_title_with_image_on_right(prs, 'Bible Reading', 'Bible', fonts['title'] - 10)
    
    # Process Bible Verses
    # The request should ideally already have the text, but if we need to fetch it:
    # We will use our bible.py logic.
    from .bible import bible_passage_auto, get_correct_copyright_message
    
    # Track used versions for copyright
    used_versions = set()

    for reading in request.bible_readings:
        verse_parts = bible_passage_auto(f"{reading.reference} ({reading.version})")
        used_versions.add(reading.version)
        
        for part in verse_parts:
             create_title_and_text_slide(f"{reading.reference} ({reading.version})", part, prs, fonts['title'], fonts['bible'])
    
    # Copyright
    for version in used_versions:
        create_title_and_text_slide("", get_correct_copyright_message(version), prs, 10, 10)

    # 5. Bulletin (Backfill slide 0)
    # Collect data
    verse_refs = [f"{r.reference} ({r.version})" for r in request.bible_readings]
    response_song_names = [s.title for s in request.response_songs]
    
    create_bulletin_slide(prs.slides[0], prs, request.date, song_names, verse_refs, response_song_names, request.speaker, request.topic, request.church_name, request.service_name)

    # 6. Response Songs
    for song in request.response_songs:
        append_song(prs, song, fonts['title'], fonts['song'], request.translate, request.language)

    # 7. Announcements & Tithing
    create_title_slide('Announcements', '', prs, fonts['title'])
    for ann in request.announcements:
        if ann.content:
            create_title_and_text_slide(ann.title, ann.content, prs, fonts['title'], fonts['song'])
        else:
            create_title_slide(ann.title, '', prs, fonts['title'])

    create_offering_slide(prs, fonts['title'], fonts['tithing'], request.offering)
    
    create_title_slide('Prayer Points', '', prs, fonts['title'])
    for point in request.prayer_points:
        create_text_slide(point, prs, fonts['song'])
    
    # 8. Mingle
    create_title_slide(request.mingle_text, '', prs, fonts['title'])

    # Output
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
